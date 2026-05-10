import base64
import io

def extract_text(pitch_text: str | None, file_base64: str | None, file_type: str | None) -> str:
    if pitch_text:
        return pitch_text.strip()

    if not file_base64 or not file_type:
        raise ValueError("Either pitch_text or file_base64+file_type must be provided")

    raw = base64.b64decode(file_base64)

    if file_type == "pdf":
        import pymupdf
        doc = pymupdf.open(stream=raw, filetype="pdf")
        return "\n\n".join(page.get_text() for page in doc).strip()

    if file_type == "pptx":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n\n".join(texts).strip()

    raise ValueError(f"Unsupported file_type: {file_type}")
